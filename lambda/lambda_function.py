import boto3
import os
import fitz  # PyMuPDF
import docx
from PIL import Image
import pytesseract
import tempfile
import traceback
from pymongo import MongoClient
from dotenv import load_dotenv
from bson import ObjectId
import pandas as pd
from pptx import Presentation
from bs4 import BeautifulSoup
import requests
import base64
import json
import time
import asyncio
from concurrent.futures import ThreadPoolExecutor


print("Starting Lambda function...")

s3 = boto3.client('s3')
# Cấu hình Tesseract binary
pytesseract.pytesseract.tesseract_cmd = "/var/task/tesseract/bin/tesseract"
os.environ["TESSDATA_PREFIX"] = "/var/task/tesseract/tesseract/share/tessdata"
os.environ["LD_LIBRARY_PATH"] = "/var/task/tesseract/lib"
api_embedding =  "https://jpwd4r42ri.execute-api.ap-southeast-2.amazonaws.com/BERT_Stage/get-embedding"
api_predict = "https://jpwd4r42ri.execute-api.ap-southeast-2.amazonaws.com/BERT_Stage/predict"

esUsername = "elastic"
esPassword = "nGa5A2umZ6Tnh1Kmeh0SeK8Y"  
esIndex = "files"  # tên index trong Elasticsearch
esEndpoint = "https://a45489d5009a463da1487c645bf718be.us-central1.gcp.cloud.es.io" 

# Tạo header auth
credentials = f"{esUsername}:{esPassword}"
authHeader = "Basic " + base64.b64encode(credentials.encode()).decode()


# --- Load .env ---
load_dotenv()
MONGO_URI = os.getenv("MONGODB_CONNECTION_STRING")

# --- Kết nối MongoDB ---
client = MongoClient(MONGO_URI)
db = client["test"]  # đổi thành tên DB của bạn
collection = db["files"]   # đổi thành tên collection chứa FileSchema


def extract_text(filepath, file_ext):
    if not os.path.exists(filepath):
        print(f"[ERROR] File not found: {filepath}")
        return ""
    if file_ext == 'txt' or file_ext == 'md':
        with open(filepath, 'r', encoding='utf-8') as f:
            return f.read()
    elif file_ext == 'pdf':
        doc = fitz.open(filepath)
        text = ""
        for page in doc:
            text += page.get_text()
        return text
    elif file_ext in ['docx', 'doc']:
        doc = docx.Document(filepath)
        return "\n".join([para.text for para in doc.paragraphs])
    elif file_ext in ['jpg', 'jpeg', 'png']:
        img = Image.open(filepath).convert('L')  # Convert to grayscale
        return pytesseract.image_to_string(img)
    elif file_ext in ['xlsx', 'csv']:
        df = pd.read_excel(filepath) if file_ext == 'xlsx' else pd.read_csv(filepath)
        return df.to_string(index=False)
    elif file_ext == "pptx":
        prs = Presentation(filepath)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        return '\n'.join(text)
    elif file_ext == "html":
        with open(filepath, 'r', encoding='utf-8') as f:
            soup = BeautifulSoup(f, 'html.parser')
            return soup.get_text(separator='\n')
    else:
        print(f"[WARNING] Unsupported file type: {file_ext}")
        return "Unsupported file type"

def chunk_text_by_chars(text: str, max_chars, overlap):
    """
    Chia text theo ký tự, giữ overlap.
    Trả về list các chunk.
    """
    if not text:
        return []
    text = text.strip()
    if len(text) <= max_chars:
        return [{
            "text": text,
            "offset_start": 0,
            "offset_end": len(text)
        }]
    chunks = []
    start = 0
    text_len = len(text)
    while start < text_len:
        end = start + max_chars
        if end >= text_len:
            chunk = text[start:text_len]
            chunks.append({
                "text": chunk,
                "offset_start": start,
                "offset_end": end
            })
            break
        # cố gắng cắt ở khoảng trắng gần cuối để không cắt giữa từ
        split_at = text.rfind(" ", start, end)
        if split_at <= start:
            split_at = end  # không tìm thấy space hợp lý -> cắt thẳng
        chunk = text[start:split_at]
        chunks.append({
                "text": chunk,
                "offset_start": start,
                "offset_end": end
            })
        start = split_at - overlap if (split_at - overlap) > start else split_at
    return chunks


def get_embedding_from_api(text: str, retries=3, delay=2):
    payload = {"text": text}
    for attempt in range(retries):
        try:
            resp = requests.post(api_embedding, json=payload, timeout=30)
            resp.raise_for_status()
            body = resp.json()
            if "embedding" in body:
                return body["embedding"][0][0]
            print(f"[WARN] Unexpected response: {list(body.keys())}")
            return None
        except Exception as e:
            print(f"[ERROR] Attempt {attempt+1} failed: {e}")
            if attempt < retries - 1:
                time.sleep(delay)
    return None

def index_doc_to_es(doc: dict, es_index: str = esIndex):
    """
    Index 1 document vào Elasticsearch cloud.
    Trả về True nếu thành công.
    """
    try:
        url = f"{esEndpoint.rstrip('/')}/{es_index}/_doc"
        headers = {
            "Content-Type": "application/json",
            "Authorization": authHeader
        }
        resp = requests.post(url, headers=headers, data=json.dumps(doc))
        if resp.status_code in (200, 201):
            return True
        else:
            print(f"[ERROR] ES index failed status {resp.status_code}: {resp.text}")
            return False
    except Exception as e:
        print(f"[ERROR] Exception during ES indexing: {e}")
        return False

def chunk_and_index(file_id, text, max_chars=1000, overlap=200):
    """
    Chia text thành chunk, gọi embedding API cho mỗi chunk, index vào ES.
    file_id: ObjectId hoặc string
    """
    fid_str = str(file_id)
    start_chunk = time.time()
    chunks = chunk_text_by_chars(text, max_chars=max_chars, overlap=overlap)
    print(f"[INFO] chunk_and_index -> file_id={fid_str}, produced {len(chunks)} chunks (max_chars={max_chars}, overlap={overlap})")
    
    def process_chunk(chunk):
        # Call embedding API
        emb = get_embedding_from_api(chunk["text"])
        if emb is None:
            print(f"[WARN] Skipping chunk due to failed embedding.")
            return

        # Build document
        doc = {
            "file_id": fid_str,
            "content": chunk["text"],
            "vector_embedding": emb,
            "offset_start": chunk["offset_start"],
            "offset_end": chunk["offset_end"]
        }

        return index_doc_to_es(doc, esIndex)

    with ThreadPoolExecutor(max_workers=8) as pool:
        results = list(pool.map(process_chunk, chunks))
    end_chunk = time.time() 
    print(f"[INFO] chunk_and_index -> file_id={fid_str}, total time: {end_chunk - start_chunk:.2f} seconds")
    print(f"[INFO] Indexed {sum(results)}/{len(chunks)} chunks for {fid_str}")


def classify(text, retries=3, delay=2):
    payload = {
        "text": text
    }
    for attempt in range(retries):
        try:
            # Gọi POST request
            response = requests.post(api_predict, json=payload)
            response.raise_for_status()
            # Parse kết quả từ model API
            result = response.json()
            predicted_class = result.get("predicted_class")
            return predicted_class
        except Exception as e:
            print(f"[ERROR] Attempt {attempt+1} failed: {e}")
            if attempt < retries - 1:
                time.sleep(delay)
    return None

# --- Predict & update MongoDB ---
def classify_and_update(file_id, text, retries=3, delay=2):
    predicted_class = classify(text, retries=retries, delay=delay)

    file_info = collection.find_one({"_id": file_id})

    print(f"File ID: {file_id}, Current Category: {file_info.get('document_category', 'N/A')}")

    # --- Cập nhật document_category ---
    result = collection.update_one(
        {"_id": file_id},  
        {"$set": {"document_category": predicted_class}}
    )
    print(f"Updated file {file_id} to category: {predicted_class} - Modified: {result.modified_count}")


# ---------- ASYNC HANDLER ----------
async def process_file(file_id, text):
    loop = asyncio.get_event_loop()
    with ThreadPoolExecutor() as pool:
        classify_future = loop.run_in_executor(pool, classify_and_update, ObjectId(file_id), text)
        index_future = loop.run_in_executor(pool, chunk_and_index, file_id, text)
        await asyncio.gather(classify_future, index_future)

def lambda_handler(event, context):
    try:
        bucket = event['Records'][0]['s3']['bucket']['name']
        key    = event['Records'][0]['s3']['object']['key']

        file_name = key.split('/')[-1]
        file_ext  = file_name.split('.')[-1].lower()

        with tempfile.NamedTemporaryFile(delete=False) as tmp_file:
            s3.download_fileobj(bucket, key, tmp_file)
            tmp_file_path = tmp_file.name

        metadata = s3.head_object(Bucket=bucket, Key=key).get('Metadata', {})
        file_id = metadata.get('file_id', None)
        text = ""
        # Extract text based on file type
        text = extract_text(tmp_file_path, file_ext)

        print(f"[INFO] File processed: {file_name}")
        print(f"[INFO] File ID: {file_id}")
        print(f"[INFO] Extracted text (truncated):\n{text[:1000]}")

        if file_id and text != "Unsupported file type" and text != "":
            asyncio.run(process_file(file_id, text))
        else:
            print(f"[WARNING] No valid file_id or text extracted for file: {file_name}")

        return {
            "statusCode": 200,
            "body": f"Successfully processed file: {file_name}"
        }

    except Exception as e:
        print(f"[ERROR] Exception occurred: {str(e)}")
        print(traceback.format_exc())
        return {
            "statusCode": 500,
            "body": f"Error processing file: {file_name if 'file_name' in locals() else 'unknown'}"
        }

    finally:
        # Clean up temp file
        if 'tmp_file_path' in locals() and os.path.exists(tmp_file_path):
            os.remove(tmp_file_path)

